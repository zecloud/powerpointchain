from langchain.chat_models import AzureChatOpenAI
from langchain.agents import Tool
from langchain.agents import AgentType, initialize_agent, load_tools
from langchain.utilities import BingSearchAPIWrapper
from langchain.callbacks import StreamlitCallbackHandler
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.callbacks import HumanApprovalCallbackHandler
from langchain.output_parsers import PydanticOutputParser
from langchain.tools import StructuredTool
from pydantic.v1 import BaseModel, Field, validator
#from msal_streamlit_authentication import msal_authentication
from st_files_connection import FilesConnection
from pptx import Presentation
from typing import List
import requests
from bs4 import BeautifulSoup
import streamlit as st
import os
from io import BytesIO 
from pathlib import Path
import json
from PIL import Image
#from langchain.utilities.dalle_image_generator import DallEAPIWrapper
import openai


openai_api_key= os.environ.get('OPENAIAPIKEY')
dalle_api_key= os.environ.get('DALLEOPENAIAPIKEY')
azure_resource_name = os.environ.get('RESOURCE_NAME',"zecloud")
dalle_resource_name = os.environ.get('DALLE_RESOURCE_NAME',"aiavatar")
azure_deployment_name = os.environ.get('DEPLOYMENT_NAME',"chatgpt")
BASE_URL = "https://"+azure_resource_name+".openai.azure.com"
DALLE_BASE_URL = "https://"+dalle_resource_name+".openai.azure.com"
API_KEY = openai_api_key
DEPLOYMENT_NAME = azure_deployment_name

# with st.sidebar:
#     token = msal_authentication(
#         auth={
#             "clientId": "bba1e505-6268-4e5b-a685-92567d288344",
#             "authority": "https://login.microsoftonline.com/9d430f15-406c-4dc1-9606-6a7406d91642",
#             "redirectUri": "/",
#             "postLogoutRedirectUri": "/"
#         },
#         cache={
#             "cacheLocation": "sessionStorage",
#             "storeAuthStateInCookie": False
#         },
#         login_request={
#             "scopes": ["files.readwrite.all"]
#         },
#         key=1)

conn = st.experimental_connection('abfs', type=FilesConnection)
#with st.sidebar:
    
        


os.environ["BING_SEARCH_URL"]=os.environ.get("BING_ENDPOINT","https://api.bing.microsoft.com//v7.0/search") 
search = BingSearchAPIWrapper()
llm =  AzureChatOpenAI(
        openai_api_base=BASE_URL,
        openai_api_version="2023-05-15",
        deployment_name=DEPLOYMENT_NAME,
        openai_api_key=API_KEY,
        openai_api_type="azure",
        streaming=True
    )

HEADERS = {
    'User-Agent': 'Mozilla/5.0 (Macintosh; Intel Mac OS X 10.15; rv:90.0) Gecko/20100101 Firefox/90.0'
}

def parse_html(content) -> str:
    soup = BeautifulSoup(content, 'html.parser')
    text_content_with_links = soup.get_text()
    return text_content_with_links

def fetch_web_page(url: str) -> str:
    response = requests.get(url, headers=HEADERS)
    return parse_html(response.content)

web_fetch_tool = Tool.from_function(
    func=fetch_web_page,
    name="WebFetcher",
    description="Useful only if you need more detail to fetch the content of a web page"
)



class Slide(BaseModel):
    title: str = Field(description="Title of the slide")
    content: List[str] = Field(description="Content of the slide")
    note: str = Field(description="Note of the slide")
    prompt: str = Field(description="Prompt of the slide")

class Slides(BaseModel):
    subject: str = Field(description="Subject of the presentation")
    slides: List[Slide] = Field(description="Slides of the presentation")
    
with st.sidebar:
    prefixfilename=st.text_input("prefixe au nom de fichier","FR_977_01_",key="prefixfilename")

def create_image(prompt):
    openai.api_base = DALLE_BASE_URL # Enter your endpoint here
    openai.api_key = dalle_api_key        # Enter your API key here
    openai.api_version = '2023-06-01-preview'
    openai.api_type = 'azure'
    
    
    # Create an image by using the image generation API
    generation_response = openai.Image.create(
        prompt=prompt,    # Enter your prompt text here
        size='1024x1024',
        n=1
    )
    imgurl = generation_response["data"][0]["url"]
    #imgurl=DallEAPIWrapper().run(prompt)
    generated_image = requests.get(imgurl).content  
    return generated_image


def save_slides(subject:str,slides: Slides):
    with conn.open("templates/template.pptx", "rb") as f:
        source_stream = BytesIO(f.read())
        prs = Presentation(source_stream)
        bullet_slide_layout = prs.slide_layouts[6]
        prs.slides[0].shapes[2].text = subject
        for i in range(0,len(slides)):
            
            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            
            body_shape = shapes.placeholders[10]

            title_shape.text = slides[i]["title"]
            tf = body_shape.text_frame
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slides[i]["note"]
            for content in slides[i]["content"]:
                p = tf.add_paragraph()
                p.text = content
                p.level = 1
            img_shape= shapes.placeholders[11]
            imgbytes = create_image(prompt=slides[i]["prompt"])
            bytes_io = BytesIO()
            bytes_io.write(imgbytes)
            img_shape.insert_picture(bytes_io)

        presentations = conn.fs.ls("powerpoints")
        stemfilename=st.session_state.get("prefixfilename","FR_977_01_")
        if(len(presentations)>0):
            idfile=int(Path(presentations[-1].split()[0]).stem[-2:])+1
            new_id = "{:02d}".format(idfile)
            stemfilename=stemfilename+new_id
        else:
            stemfilename=stemfilename+"00"
        filenamefinal=stemfilename+ " "+subject[:50]+".pptx"
        target_stream = BytesIO()
        prs.save(target_stream)
        with conn.open("powerpoints/"+filenamefinal, "wb") as f:
            f.write(target_stream.getvalue())
        return filenamefinal

tool_save_slide = StructuredTool.from_function(save_slides,description="Save the slides to a local file",args_schema=Slides)

parser = PydanticOutputParser(pydantic_object=Slides)

prompt = PromptTemplate(
    input_variables=["query"],
    partial_variables={"format_instructions": parser.get_format_instructions()},
    template="""You are a software developer creating an automated PowerPoint generation program.
         You decide the content that goes into each slide of the PowerPoint.
         Each slide typically consists of a topic, introduction, main points, conclusion, and references.
         Each slide also contains detailed notes about the topic that add content to present the information.
         Follow the rules below:\n
         Use as much slides as you need.\n
         Each slide contains 'title', 'content' 'note' and 'prompt'.\n
         'content' is a bullet point that breaks down the core content into brief, step-by-step chunks.\n
         All of the slides contain images.\n
         Write a two sentence description to create a prompt to will accurately create an image using DALL-E and save it into 'prompt'.\n
         {format_instructions}\n
         {query}\n"""
        # Summarize and extract the key contents from the user's input around 5 slides.\n
         #Focus on nouns and adjectives and separate them with commas so that 'prompt' is a good visual representation of 'content'.\n
         #Set the above contents as keys named 'title', 'content', and 'prompt'.\n
         #Save the results of each slide as a JSON list.\n
         #Output the final output in JSON format.\n
         #Make sure output JSON can be parsed with Python `json.loads()`.\n
         #Must return JSON format only\n
)

llm_chain = LLMChain(llm=llm, prompt=prompt)

# initialize the LLM tool
llm_tool = Tool(
    name='Slide Generator',
    func=llm_chain.run,
    description='useful to generate slides for your next presentation'
)


searchurls=StructuredTool.from_function(search.results,name= "SearchMetadata",description="use it for search to get the urls to fetch pages")

#tools = load_tools(["llm-math"], llm=llm)
tools=[Tool(
        name = "Search",
        func=search.run,
        description="use it for search and useful for when you need to answer questions about current events or the current state of the world",
        
    ),searchurls,llm_tool,tool_save_slide,web_fetch_tool]
agent = initialize_agent(
    tools, llm, agent=AgentType.STRUCTURED_CHAT_ZERO_SHOT_REACT_DESCRIPTION, verbose=True
)

if prompt := st.chat_input():
    st.chat_message("user").write(prompt)
    with st.chat_message("assistant"):
        st_callback = StreamlitCallbackHandler(st.container())
        response = agent.run(prompt, callbacks=[st_callback])
        st.write(response)